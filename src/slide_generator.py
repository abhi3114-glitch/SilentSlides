"""
Slide Generator Module
Supports PDF, PPTX, and Markdown export
"""

from typing import List, Dict
from pathlib import Path
import json
from datetime import datetime

# PDF generation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from config import THEMES_DIR, OUTPUT_DIR, DEFAULT_THEME


class Theme:
    """Theme configuration loader"""
    
    def __init__(self, theme_name: str = DEFAULT_THEME):
        self.name = theme_name
        self.config = self._load_theme(theme_name)
    
    def _load_theme(self, theme_name: str) -> Dict:
        """Load theme from JSON file"""
        theme_path = THEMES_DIR / f"{theme_name}.json"
        
        if not theme_path.exists():
            # Return default theme
            return self._get_default_theme()
        
        with open(theme_path, 'r') as f:
            return json.load(f)
    
    def _get_default_theme(self) -> Dict:
        """Default clean theme"""
        return {
            "name": "clean",
            "colors": {
                "background": "#FFFFFF",
                "title": "#1a1a1a",
                "text": "#333333",
                "accent": "#0066cc"
            },
            "fonts": {
                "title": "Helvetica-Bold",
                "heading": "Helvetica-Bold",
                "body": "Helvetica"
            },
            "sizes": {
                "title": 28,
                "heading": 24,
                "body": 14
            }
        }


class PDFGenerator:
    """Generate PDF slides using ReportLab"""
    
    def __init__(self, theme: Theme):
        self.theme = theme
    
    def generate(self, data: Dict, output_path: str) -> str:
        """
        Generate PDF presentation
        
        Args:
            data: Processed slide data
            output_path: Output file path
        
        Returns:
            Path to generated file
        """
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )
        
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=32,
            textColor=self.theme.config['colors']['title'],
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=self.theme.config['colors']['title'],
            spaceAfter=20
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=14,
            textColor=self.theme.config['colors']['text'],
            spaceAfter=10,
            leftIndent=20
        )
        
        # Title slide
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph("<b>AI-Generated Slide Deck</b>", title_style))
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", body_style))
        story.append(PageBreak())
        
        # Topic slides
        for topic in data.get('topics', []):
            story.append(Spacer(1, 0.5*inch))
            story.append(Paragraph(f"<b>{topic['title']}</b>", heading_style))
            story.append(Spacer(1, 0.3*inch))
            
            for bullet in topic['bullets']:
                story.append(Paragraph(f"• {bullet}", body_style))
            
            story.append(PageBreak())
        
        # Summary slide
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph("<b>Summary</b>", title_style))
        story.append(Spacer(1, 0.3*inch))
        
        summary_text = f"Generated {len(data.get('topics', []))} topics from {data.get('total_sentences', 0)} sentences"
        story.append(Paragraph(summary_text, body_style))
        
        # Build PDF
        doc.build(story)
        
        return output_path


class PPTXGenerator:
    """Generate PowerPoint slides using python-pptx"""
    
    def __init__(self, theme: Theme):
        self.theme = theme
    
    def generate(self, data: Dict, output_path: str) -> str:
        """
        Generate PPTX presentation
        
        Args:
            data: Processed slide data
            output_path: Output file path
        
        Returns:
            Path to generated file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "AI-Generated Slide Deck"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Topic slides
        for topic in data.get('topics', []):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = topic['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            
            # Bullets
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(5)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(topic['bullets']):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = bullet
                p.font.size = Pt(18)
                p.level = 0
                p.space_after = Pt(12)
        
        # Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        summary_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5), Inches(6), Inches(1)
        )
        summary_frame = summary_box.text_frame
        summary_text = f"Generated {len(data.get('topics', []))} topics from {data.get('total_sentences', 0)} sentences"
        summary_frame.text = summary_text
        summary_para = summary_frame.paragraphs[0]
        summary_para.font.size = Pt(20)
        summary_para.alignment = PP_ALIGN.CENTER
        
        # Save
        prs.save(output_path)
        
        return output_path


class MarkdownGenerator:
    """Generate Markdown slides"""
    
    def __init__(self, theme: Theme):
        self.theme = theme
    
    def generate(self, data: Dict, output_path: str) -> str:
        """
        Generate Markdown presentation
        
        Args:
            data: Processed slide data
            output_path: Output file path
        
        Returns:
            Path to generated file
        """
        lines = []
        
        # Title slide
        lines.append("# AI-Generated Slide Deck")
        lines.append("")
        lines.append(f"*Generated on {datetime.now().strftime('%B %d, %Y')}*")
        lines.append("")
        lines.append("---")
        lines.append("")
        
        # Topic slides
        for topic in data.get('topics', []):
            lines.append(f"## {topic['title']}")
            lines.append("")
            
            for bullet in topic['bullets']:
                lines.append(f"- {bullet}")
            
            lines.append("")
            lines.append("---")
            lines.append("")
        
        # Summary slide
        lines.append("## Summary")
        lines.append("")
        summary_text = f"Generated {len(data.get('topics', []))} topics from {data.get('total_sentences', 0)} sentences"
        lines.append(summary_text)
        
        # Write file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        
        return output_path


class SlideGenerator:
    """Main slide generator with multi-format support"""
    
    def __init__(self, theme_name: str = DEFAULT_THEME):
        self.theme = Theme(theme_name)
        self.pdf_gen = PDFGenerator(self.theme)
        self.pptx_gen = PPTXGenerator(self.theme)
        self.md_gen = MarkdownGenerator(self.theme)
    
    def generate_all(self, data: Dict, base_name: str = "slides") -> Dict[str, str]:
        """
        Generate slides in all formats
        
        Args:
            data: Processed slide data
            base_name: Base filename (without extension)
        
        Returns:
            Dict mapping format -> file path
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{base_name}_{timestamp}"
        
        outputs = {}
        
        # PDF
        pdf_path = OUTPUT_DIR / f"{filename}.pdf"
        outputs['pdf'] = self.pdf_gen.generate(data, str(pdf_path))
        
        # PPTX
        pptx_path = OUTPUT_DIR / f"{filename}.pptx"
        outputs['pptx'] = self.pptx_gen.generate(data, str(pptx_path))
        
        # Markdown
        md_path = OUTPUT_DIR / f"{filename}.md"
        outputs['md'] = self.md_gen.generate(data, str(md_path))
        
        return outputs
